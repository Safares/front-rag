#!/usr/bin/env python3
"""
Worker chamado pelo server.js.
Uso:
  python rag_worker.py upload --provider openai|gemini --key KEY --files-file PATH_JSON
  python rag_worker.py query  --provider openai|gemini --key KEY --store STORE_ID --question "..."
"""

import sys
import os
os.environ.setdefault('PYTHONIOENCODING', 'utf-8')
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')
    sys.stderr.reconfigure(encoding='utf-8')
import json
import argparse
import time
import re
import unicodedata
from pathlib import Path


def _emit(prefix: str, payload) -> None:
    line = payload if isinstance(payload, str) else json.dumps(payload, ensure_ascii=False)
    print(f"{prefix}{line}", flush=True)


def progress(msg: str) -> None:
    _emit("PROGRESS:", msg)


def result(data: dict) -> None:
    import sys
    sys.stderr.flush()
    sys.stdout.flush()
    _emit("RESULT:", data)
    sys.stdout.flush()


def error(msg: str) -> None:
    _emit("ERROR:", msg)
    sys.exit(1)


# ─── Leitura de arquivo ───────────────────────────────────────

def _extract_pdf_text(path: Path) -> str:
    """Extrai texto de PDF. Usa OCR como fallback para PDFs escaneados."""
    parts = []
    num_pages = 1

    # Tentativa 1: pdfplumber — guarda num_pages para evitar segunda abertura
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            num_pages = max(len(pdf.pages), 1)
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    parts.append(t)
    except ImportError:
        pass
    except Exception:
        parts = []

    # Tentativa 2: pypdf (fallback se pdfplumber falhou ou não está instalado)
    if not parts:
        try:
            import pypdf
            reader = pypdf.PdfReader(str(path))
            num_pages = max(len(reader.pages), 1)
            for page in reader.pages:
                t = page.extract_text() or ""
                if t.strip():
                    parts.append(t)
        except ImportError:
            pass
        except Exception:
            parts = []

    combined = "\n\n".join(parts)

    # < 50 chars por página em média → provável PDF escaneado
    is_scanned = len(combined.strip()) < 50 * num_pages

    if is_scanned:
        progress(f"PDF escaneado detectado em {path.name}, usando OCR (pode demorar)...")
        try:
            from pdf2image import convert_from_path
            import pytesseract
            images = convert_from_path(str(path), dpi=200)
            ocr_parts = []
            for i, img in enumerate(images, 1):
                progress(f"  OCR página {i}/{len(images)}...")
                text = pytesseract.image_to_string(img, lang="por+eng")
                if text.strip():
                    ocr_parts.append(text)
            if ocr_parts:
                return "\n\n".join(ocr_parts)
            progress(f"OCR não produziu texto para {path.name}.")
        except ImportError:
            progress("pytesseract/pdf2image não instalados. Instale com: pip install pytesseract pdf2image")
        except Exception as e:
            progress(f"Erro no OCR: {e}")

    # Retorna string vazia se nenhuma extração funcionou — o caller marcará como falha
    return combined


def read_file_as_bytes(path: Path) -> tuple[bytes, str]:
    """Retorna (conteudo_bytes, nome_para_upload)."""
    ext = path.suffix.lower()

    if ext in (".txt", ".md"):
        return path.read_bytes(), path.name

    if ext == ".pdf":
        text = _extract_pdf_text(path)
        return text.encode("utf-8"), path.stem + ".txt"

    if ext in (".xlsx", ".xls"):
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        lines = []
        for ws in wb.worksheets:
            lines.append(f"[Planilha: {ws.title}]")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    lines.append(" | ".join(cells))
        return "\n".join(lines).encode("utf-8"), path.stem + ".txt"

    if ext == ".csv":
        import csv
        lines = []
        for enc in ("utf-8-sig", "utf-8", "latin-1"):
            try:
                with open(path, newline="", encoding=enc) as f:
                    reader = csv.reader(f)
                    for row in reader:
                        cells = [str(c) for c in row if c is not None]
                        if cells:
                            lines.append(" | ".join(cells))
                break
            except UnicodeDecodeError:
                continue
        return "\n".join(lines).encode("utf-8"), path.stem + ".txt"

    if ext == ".docx":
        try:
            from docx import Document
            doc = Document(str(path))
            parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text)
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
            return "\n\n".join(parts).encode("utf-8"), path.stem + ".txt"
        except ImportError:
            error("Instale python-docx: pip install python-docx")

    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                slide_parts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_parts.append(text)
                if slide_parts:
                    parts.append(f"## Slide {i}: {slide_parts[0]}\n\n" + "\n".join(slide_parts))
            return "\n\n".join(parts).encode("utf-8"), path.stem + ".txt"
        except ImportError:
            error("Instale python-pptx: pip install python-pptx")

    error(f"Formato nao suportado: {ext}")


# ─── Helpers Gemini ───────────────────────────────────────────

def _ascii_safe(text: str) -> str:
    return unicodedata.normalize("NFKD", text).encode("ascii", "ignore").decode("ascii")


def _ascii_filename(name: str) -> str:
    if "." in name:
        stem, ext = name.rsplit(".", 1)
        safe = re.sub(r"[^A-Za-z0-9._-]", "_", _ascii_safe(stem)).strip("_") or "upload"
        return f"{safe}.{ext}"
    safe = re.sub(r"[^A-Za-z0-9._-]", "_", _ascii_safe(name)).strip("_") or "upload"
    return safe


# ─── Chunking semântico ───────────────────────────────────────

CHUNK_MAX_CHARS = 7000
CHUNK_OVERLAP_CHARS = 700


def chunk_text(text: str) -> list[str]:
    """Break text into overlapping chunks at paragraph/section boundaries."""
    if len(text) <= CHUNK_MAX_CHARS:
        return [text]

    chunks = []
    lines = text.split('\n')
    current_lines = []
    current_len = 0

    for line in lines:
        line_len = len(line) + 1
        is_section_break = line.startswith('#') or (not line.strip() and current_len > CHUNK_MAX_CHARS // 2)

        if current_len + line_len > CHUNK_MAX_CHARS and current_lines and is_section_break:
            chunks.append('\n'.join(current_lines))
            # Overlap: keep last CHUNK_OVERLAP_CHARS worth of lines
            overlap = []
            overlap_len = 0
            for prev in reversed(current_lines):
                if overlap_len + len(prev) + 1 > CHUNK_OVERLAP_CHARS:
                    break
                overlap.insert(0, prev)
                overlap_len += len(prev) + 1
            current_lines = overlap
            current_len = overlap_len
        elif current_len + line_len > CHUNK_MAX_CHARS * 1.5 and current_lines:
            # Force break even without section boundary if too large
            chunks.append('\n'.join(current_lines))
            overlap = []
            overlap_len = 0
            for prev in reversed(current_lines):
                if overlap_len + len(prev) + 1 > CHUNK_OVERLAP_CHARS:
                    break
                overlap.insert(0, prev)
                overlap_len += len(prev) + 1
            current_lines = overlap
            current_len = overlap_len

        current_lines.append(line)
        current_len += line_len

    if current_lines:
        chunks.append('\n'.join(current_lines))

    return chunks


# ─── Verificação de chunking via LLM ──────────────────────────

CHUNK_CHECK_SAMPLE_CHARS = 15000


def check_already_chunked(provider: str, api_key: str, text: str) -> dict:
    """Pergunta à LLM se o texto já está bem dividido em chunks/seções,
    para decidir se o chunk_text() automático deve ser pulado."""
    if len(text) <= CHUNK_MAX_CHARS:
        return {"already_chunked": True, "note": "Texto cabe em um único chunk (menos de 7000 caracteres)."}

    sample = text[:CHUNK_CHECK_SAMPLE_CHARS]
    prompt = (
        "Você é um especialista em preparar textos para RAG (Retrieval-Augmented Generation).\n\n"
        "Analise se o texto abaixo JÁ ESTÁ bem estruturado em chunks/seções (ex: cabeçalhos markdown, "
        "separadores claros como \"---\", numeração de seções, blocos coerentes de tamanho razoável), "
        "de forma que dividir automaticamente a cada ~7000 caracteres NÃO cortaria uma ideia no meio.\n\n"
        "Responda EXATAMENTE neste formato, sem mais nada:\n"
        "VEREDITO: SIM ou NAO\n"
        "MOTIVO: <uma frase curta>\n\n"
        "--- TEXTO (pode estar truncado) ---\n"
        f"{sample}"
    )

    try:
        if provider == "openai":
            from openai import OpenAI
            client = OpenAI(api_key=api_key)
            resp = client.responses.create(model="gpt-4o-mini", input=prompt)
            raw = resp.output_text.strip()
        else:
            from google import genai
            client = genai.Client(api_key=api_key)
            resp = client.models.generate_content(model="gemini-flash-latest", contents=prompt)
            raw = (resp.text or "").strip()

        verdict_match = re.search(r"VEREDITO:\s*(SIM|NAO|N[ÃA]O)", raw, re.IGNORECASE)
        motivo_match  = re.search(r"MOTIVO:\s*(.+)", raw, re.IGNORECASE)
        already = bool(verdict_match) and verdict_match.group(1).upper().startswith("SIM")
        note = motivo_match.group(1).strip() if motivo_match else raw[:200]
        return {"already_chunked": already, "note": note}
    except Exception as e:
        return {"already_chunked": False, "note": f"Não foi possível verificar automaticamente ({e}); será dividido pelo método padrão."}


# ─── OpenAI ───────────────────────────────────────────────────

def upload_openai(api_key: str, file_entries: list) -> dict:
    import io
    from openai import OpenAI

    client = OpenAI(api_key=api_key)
    total      = len(file_entries)
    first_stem = Path(file_entries[0]['name']).stem if file_entries else 'kb'

    progress("Criando Vector Store...")
    vs = client.vector_stores.create(name=f"kb-{first_stem}")
    progress(f"Vector Store criado: {vs.id}")

    uploaded: list[str] = []
    failed:   list[str] = []

    for i, entry in enumerate(file_entries, 1):
        fp   = Path(entry['path'])
        name = entry['name']
        try:
            content_bytes, upload_name = read_file_as_bytes(fp)
            text = content_bytes.decode("utf-8", errors="replace")
            chunks = chunk_text(text)
            stem = Path(upload_name).stem

            if len(chunks) > 1:
                progress(f"[{i}/{total}] {name}: dividido em {len(chunks)} chunks.")

            chunk_ok = False
            for n, chunk in enumerate(chunks, 1):
                chunk_name = f"{stem}_parte{n}.txt" if len(chunks) > 1 else upload_name
                chunk_bytes = chunk.encode("utf-8")
                progress(f"[{i}/{total}] Enviando {chunk_name} ({len(chunk_bytes) / 1024:.1f} KB)...")
                client.vector_stores.files.upload_and_poll(
                    vector_store_id=vs.id,
                    file=(chunk_name, io.BytesIO(chunk_bytes), "text/plain"),
                )
                progress(f"[{i}/{total}] {chunk_name} indexado.")
                chunk_ok = True

            if chunk_ok:
                uploaded.append(name)
        except Exception as e:
            failed.append(name)
            progress(f"[{i}/{total}] ERRO em {name}: {e}")

    if not uploaded:
        error("Nenhum arquivo pôde ser indexado.")

    progress("Indexacao concluida!")
    return {"store_id": vs.id, "store_name": vs.name, "provider": "openai",
            "files_uploaded": uploaded, "files_failed": failed}


def query_openai(api_key: str, store_id: str, question: str) -> dict:
    from openai import OpenAI
    client = OpenAI(api_key=api_key)
    
    try:
        progress(f"Consultando OpenAI com store: {store_id}")
        response = client.beta.threads.runs.submit_tool_outputs(
            model="gpt-4o",
            tools=[{"type": "file_search", "vector_store_ids": [store_id]}],
            messages=[{"role": "user", "content": question}]
        )
        text = ""
        if hasattr(response, 'content'):
            for item in response.content:
                if hasattr(item, 'text'):
                    text += item.text
        progress(f"Resposta OpenAI recebida: {len(text)} caracteres")
    except Exception as e:
        progress(f"Erro ao usar API beta ({e}). Tentando chat completions...")
        try:
            response = client.chat.completions.create(
                model="gpt-4o",
                messages=[{"role": "user", "content": question}],
            )
            text = response.choices[0].message.content if response.choices else ""
            progress(f"Resposta OpenAI (fallback): {len(text)} caracteres")
        except Exception as e2:
            progress(f"Erro no fallback também: {e2}")
            text = f"[Erro ao gerar resposta: {str(e2)[:100]}]"

    if not text or text.strip() == "":
        text = "[Nenhuma resposta encontrada nos documentos]"

    # Extrair citations das annotations
    citations = []
    seen = set()
    if hasattr(response, 'content'):
        for item in response.content:
            annotations = getattr(item, 'annotations', None) or []
            for ann in annotations:
                if getattr(ann, 'type', None) == 'file_citation':
                    fname = getattr(ann, 'filename', None) or getattr(ann, 'file_id', 'desconhecido')
                    if fname not in seen:
                        seen.add(fname)
                        citations.append(fname)

    return {"answer": text.strip(), "citations": citations}


# ─── Gemini ───────────────────────────────────────────────────

def upload_gemini(api_key: str, file_entries: list) -> dict:
    import tempfile
    import os
    from google import genai

    client     = genai.Client(api_key=api_key)
    total      = len(file_entries)
    first_stem = Path(file_entries[0]['name']).stem if file_entries else 'kb'

    progress("Criando File Search Store...")
    store = client.file_search_stores.create(
        config={"display_name": _ascii_safe(f"kb-{first_stem}")}
    )
    progress(f"Store criado: {store.name}")

    uploaded: list[str] = []
    failed:   list[str] = []

    for i, entry in enumerate(file_entries, 1):
        fp   = Path(entry['path'])
        name = entry['name']
        try:
            content_bytes, upload_name = read_file_as_bytes(fp)
            text = content_bytes.decode("utf-8", errors="replace")
            chunks = chunk_text(text)
            stem = Path(upload_name).stem

            if len(chunks) > 1:
                progress(f"[{i}/{total}] {name}: dividido em {len(chunks)} chunks.")

            chunk_ok = False
            for n, chunk in enumerate(chunks, 1):
                chunk_name = _ascii_filename(f"{stem}_parte{n}.txt" if len(chunks) > 1 else upload_name)
                chunk_bytes = chunk.encode("utf-8")
                progress(f"[{i}/{total}] Enviando {chunk_name} ({len(chunk_bytes) / 1024:.1f} KB)...")
                tmp = tempfile.NamedTemporaryFile(suffix=".txt", delete=False)
                try:
                    tmp.write(chunk_bytes)
                    tmp.close()
                    operation = client.file_search_stores.upload_to_file_search_store(
                        file=tmp.name,
                        file_search_store_name=store.name,
                        config={"display_name": _ascii_safe(chunk_name)},
                    )
                finally:
                    os.unlink(tmp.name)
                progress(f"[{i}/{total}] Indexando {chunk_name}...")
                while not operation.done:
                    time.sleep(5)
                    operation = client.operations.get(operation)
                    progress(f"[{i}/{total}] ...aguardando {chunk_name}")
                progress(f"[{i}/{total}] {chunk_name} indexado.")
                chunk_ok = True

            if chunk_ok:
                uploaded.append(name)
        except Exception as e:
            failed.append(name)
            progress(f"[{i}/{total}] ERRO em {name}: {e}")

    if not uploaded:
        error("Nenhum arquivo pôde ser indexado.")

    progress("Indexacao concluida!")
    return {"store_id": store.name, "store_name": store.name, "provider": "gemini",
            "files_uploaded": uploaded, "files_failed": failed}


def query_gemini(api_key: str, store_name: str, question: str) -> dict:
    from google import genai
    from google.genai import types

    client = genai.Client(api_key=api_key)
    
    prompt_with_citation = f"{question}\n\nAo final da resposta, liste as fontes usadas no formato:\n[Fontes: nome_do_arquivo_1, nome_do_arquivo_2]"
    
    text = ""
    progress(f"[DEBUG] Iniciando query_gemini com store_name={store_name}, question={question[:50]}...")
    
    try:
        progress(f"[DEBUG] Tentando com file search...")
        progress(f"[DEBUG] Store name being used: {store_name}")
        response = client.models.generate_content(
            model="gemini-flash-latest",
            contents=prompt_with_citation,
            config=types.GenerateContentConfig(
                tools=[
                    types.Tool(
                        file_search=types.FileSearch(
                            file_search_store_names=[store_name]
                        )
                    )
                ]
            ),
        )
        
        progress(f"[DEBUG] Response object: {response}")
        progress(f"[DEBUG] Response.text exists: {hasattr(response, 'text')}")
        if hasattr(response, 'text'):
            progress(f"[DEBUG] Response.text value: {repr(response.text)}")
        
        text = response.text if hasattr(response, 'text') and response.text else ""
        progress(f"[DEBUG] Resposta obtida: {len(text)} chars, conteúdo: {repr(text[:100])}")
    except Exception as e:
        progress(f"[DEBUG] Erro com file search: {type(e).__name__}: {str(e)}")
        import traceback
        progress(f"[DEBUG] Traceback: {traceback.format_exc()}")
        try:
            progress(f"[DEBUG] Tentando fallback sem file search...")
            response = client.models.generate_content(
                model="gemini-flash-latest",
                contents=question
            )
            text = response.text if hasattr(response, 'text') and response.text else ""
            progress(f"[DEBUG] Resposta fallback: {len(text)} chars")
        except Exception as e2:
            progress(f"[DEBUG] Erro fallback: {type(e2).__name__}: {str(e2)}")
            text = f"[Erro: {str(e2)[:50]}]"
    
    if not text or text.strip() == "":
        text = "[Nenhuma resposta - tente reformular a pergunta]"
        progress(f"[DEBUG] Resposta estava vazia, usando placeholder")
    
    progress(f"[DEBUG] Text final: {len(text)} chars")
    
    # Extrair citations
    citations = []
    match = re.search(r'\[Fontes?:\s*([^\]]+)\]', text, re.IGNORECASE)
    if match:
        raw = match.group(1)
        citations = [c.strip() for c in raw.split(',') if c.strip()]
        text = text[:match.start()].rstrip()
    
    progress(f"[DEBUG] Retornando resposta com {len(citations)} citations")
    return {"answer": text.strip(), "citations": citations}


# ─── Web crawl ────────────────────────────────────────────────

MAX_PAGES = 300


def _links_from_html(html: str, base_url: str, base_domain: str) -> list[str]:
    from bs4 import BeautifulSoup
    from urllib.parse import urljoin, urlparse, urldefrag
    soup = BeautifulSoup(html, "lxml")
    found = []
    for a in soup.find_all("a", href=True):
        href, _ = urldefrag(urljoin(base_url, a["href"]))
        p = urlparse(href)
        if p.netloc == base_domain and p.scheme in ("http", "https"):
            found.append(href)
    return found


def crawl_site(url: str, depth: int, use_js: bool = False) -> list[dict]:
    """Descobre páginas dentro do mesmo domínio.
    depth=0 significa sem limite (até MAX_PAGES).
    use_js=True usa Playwright para renderizar JavaScript.
    """
    from urllib.parse import urlparse, urldefrag
    visited: set[str] = set()
    pages:   list[dict] = []
    base_domain = urlparse(url).netloc
    effective_depth = depth if depth > 0 else 999

    if use_js:
        _crawl_js(url, effective_depth, visited, pages, base_domain)
    else:
        _crawl_requests(url, effective_depth, 1, visited, pages, base_domain)

    return pages


def _crawl_requests(current: str, depth: int, level: int,
                    visited: set, pages: list, base_domain: str) -> None:
    import requests
    from urllib.parse import urldefrag

    if len(pages) >= MAX_PAGES:
        return
    current, _ = urldefrag(current)
    if current in visited:
        return
    visited.add(current)

    headers = {"User-Agent": "Mozilla/5.0 (RAG-crawler/1.0)"}
    try:
        r = requests.get(current, timeout=12, headers=headers)
        if "text/html" not in r.headers.get("Content-Type", ""):
            return
        from bs4 import BeautifulSoup
        soup  = BeautifulSoup(r.content, "lxml")
        title = (soup.title.string or current).strip()
        pages.append({"title": title, "url": current})
        progress(f"[{len(pages)}] {title}")

        if level < depth:
            for href in _links_from_html(r.text, current, base_domain):
                _crawl_requests(href, depth, level + 1, visited, pages, base_domain)
    except Exception as e:
        progress(f"Ignorado ({current}): {e}")


def _crawl_js(start: str, depth: int,
              visited: set, pages: list, base_domain: str) -> None:
    """Crawl usando Playwright (para SPAs / JavaScript)."""
    from urllib.parse import urldefrag
    from playwright.sync_api import sync_playwright

    queue = [(start, 1)]

    with sync_playwright() as pw:
        browser = pw.chromium.launch(headless=True)
        page    = browser.new_page()
        page.set_extra_http_headers({"User-Agent": "Mozilla/5.0 (RAG-crawler/1.0)"})

        while queue and len(pages) < MAX_PAGES:
            current, level = queue.pop(0)
            current, _ = urldefrag(current)
            if current in visited:
                continue
            visited.add(current)

            try:
                page.goto(current, wait_until="networkidle", timeout=25000)
                title = page.title() or current
                html  = page.content()
                pages.append({"title": title.strip(), "url": current})
                progress(f"[{len(pages)}] {title.strip()}")

                if level < depth:
                    for href in _links_from_html(html, current, base_domain):
                        if href not in visited:
                            queue.append((href, level + 1))
            except Exception as e:
                progress(f"Ignorado ({current}): {e}")

        browser.close()


# ─── Scrape URLs → RAG ────────────────────────────────────────

_JUNK_KEYWORDS = [
    "cookie", "banner", "breadcrumb", "sidebar", "popup", "modal",
    "newsletter", "advertisement", "ads", "social", "share",
    "comment", "related",
]


def _extract_text(html: bytes, url: str) -> str:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(html, "lxml")

    # Remove tags estruturais de navegação/lixo
    for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
        tag.decompose()

    # Remove elementos por classes/IDs de lixo comuns
    for kw in _JUNK_KEYWORDS:
        for el in soup.select(f'[class*="{kw}"], [id*="{kw}"]'):
            el.decompose()

    # Converte tabelas para Markdown antes de extrair texto
    for table in soup.find_all("table"):
        rows = []
        for tr in table.find_all("tr"):
            cells = [td.get_text(strip=True) for td in tr.find_all(["td", "th"])]
            rows.append("| " + " | ".join(cells) + " |")
        if rows:
            # Insere separador após header (primeira linha)
            first_row = table.find("tr")
            col_count = len(first_row.find_all(["td", "th"])) if first_row else 0
            separator = "| " + " | ".join(["---"] * col_count) + " |"
            rows.insert(1, separator)
            md_table = "\n".join(rows)
            table.replace_with(soup.new_string(f"\n\n{md_table}\n\n"))

    # Marca blocos de código com backticks
    for pre in soup.find_all("pre"):
        code_text = pre.get_text()
        if not code_text.strip():
            continue
        pre.replace_with(soup.new_string(f"\n```\n{code_text}\n```\n"))

    for code in soup.find_all("code"):
        if not code.find_parent("pre"):  # apenas code inline (pre já foi tratado)
            code_text = code.get_text()
            code.replace_with(soup.new_string(f"`{code_text}`"))

    title = (soup.title.string if soup.title else None) or url
    title = title.strip()
    body  = soup.get_text(separator="\n", strip=True)
    return f"## {title}\n\nURL: {url}\n\n{body}"


def scrape_and_upload(api_key: str, provider: str, urls: list[str], name: str) -> dict:
    import io
    import requests

    headers = {"User-Agent": "Mozilla/5.0 (RAG-crawler/1.0)"}
    parts   = []

    for i, url in enumerate(urls, 1):
        try:
            r = requests.get(url, timeout=15, headers=headers)
            parts.append(_extract_text(r.content, url))
            progress(f"[{i}/{len(urls)}] Extraído: {url}")
        except Exception as e:
            progress(f"[{i}/{len(urls)}] Erro em {url}: {e}")

    if not parts:
        error("Nenhuma página pôde ser extraída.")

    combined    = "\n\n---\n\n".join(parts)
    upload_stem = re.sub(r"[^A-Za-z0-9_-]", "_", name)[:60]
    upload_name = upload_stem + ".txt"
    chunks      = chunk_text(combined)

    if len(chunks) > 1:
        progress(f"Texto combinado dividido em {len(chunks)} chunks para melhor precisão RAG.")

    if provider == "openai":
        from openai import OpenAI
        client = OpenAI(api_key=api_key)

        progress("Criando Vector Store...")
        vs = client.vector_stores.create(name=f"kb-{name[:40]}")
        progress(f"Vector Store criado: {vs.id}")

        for n, chunk in enumerate(chunks, 1):
            chunk_name  = f"{upload_stem}_parte{n}.txt" if len(chunks) > 1 else upload_name
            chunk_bytes = chunk.encode("utf-8")
            progress(f"Enviando {chunk_name} ({len(chunk_bytes) / 1024:.1f} KB)...")
            client.vector_stores.files.upload_and_poll(
                vector_store_id=vs.id,
                file=(chunk_name, io.BytesIO(chunk_bytes), "text/plain"),
            )
            progress(f"{chunk_name} indexado.")

        progress("Indexacao concluida!")
        return {"store_id": vs.id, "store_name": vs.name, "provider": "openai"}

    else:
        from google import genai
        client = genai.Client(api_key=api_key)
        import tempfile, os

        progress("Criando File Search Store...")
        store = client.file_search_stores.create(
            config={"display_name": _ascii_safe(f"kb-{name[:40]}")}
        )
        progress(f"Store criado: {store.name}")

        for n, chunk in enumerate(chunks, 1):
            chunk_name  = f"{upload_stem}_parte{n}.txt" if len(chunks) > 1 else upload_name
            chunk_bytes = chunk.encode("utf-8")
            tmp = tempfile.NamedTemporaryFile(suffix=".txt", delete=False)
            try:
                tmp.write(chunk_bytes)
                tmp.close()

                progress(f"Enviando {chunk_name} ({len(chunk_bytes) / 1024:.1f} KB)...")
                operation = client.file_search_stores.upload_to_file_search_store(
                    file=tmp.name,
                    file_search_store_name=store.name,
                    config={"display_name": _ascii_safe(chunk_name)},
                )
            finally:
                os.unlink(tmp.name)

            progress(f"Indexando {chunk_name} (pode levar alguns minutos)...")
            while not operation.done:
                time.sleep(5)
                operation = client.operations.get(operation)
                progress(f"...aguardando {chunk_name}")
            progress(f"{chunk_name} indexado.")

        progress("Indexacao concluida!")
        return {"store_id": store.name, "store_name": store.name, "provider": "gemini"}


# ─── Upload por texto (extract→confirm flow) ──────────────────

def _upload_texts_openai(api_key: str, texts: list, name: str) -> dict:
    import io
    from openai import OpenAI
    client = OpenAI(api_key=api_key)

    progress("Criando Vector Store...")
    vs = client.vector_stores.create(name=f"kb-{name[:40]}")
    progress(f"Vector Store criado: {vs.id}")

    uploaded, failed = [], []
    total = len(texts)

    for i, entry in enumerate(texts, 1):
        fname = entry['name']
        text  = entry['text']
        try:
            if entry.get('already_chunked'):
                chunks = [text]
                progress(f"[{i}/{total}] {fname}: já está bem dividido, enviado como está (sem re-cortar).")
            else:
                chunks = chunk_text(text)
                if len(chunks) > 1:
                    progress(f"[{i}/{total}] {fname}: dividido em {len(chunks)} chunks.")
            stem   = Path(fname).stem
            for n, chunk in enumerate(chunks, 1):
                cname = f"{stem}_parte{n}.txt" if len(chunks) > 1 else f"{stem}.txt"
                cbytes = chunk.encode("utf-8")
                progress(f"[{i}/{total}] Enviando {cname} ({len(cbytes)/1024:.1f} KB)...")
                client.vector_stores.files.upload_and_poll(
                    vector_store_id=vs.id,
                    file=(cname, io.BytesIO(cbytes), "text/plain"),
                )
                progress(f"[{i}/{total}] {cname} indexado.")
            uploaded.append(fname)
        except Exception as e:
            failed.append(fname)
            progress(f"[{i}/{total}] ERRO em {fname}: {e}")

    if not uploaded:
        error("Nenhum texto pôde ser indexado.")
    progress("Indexacao concluida!")
    return {"store_id": vs.id, "store_name": vs.name, "provider": "openai",
            "files_uploaded": uploaded, "files_failed": failed}


def _upload_texts_gemini(api_key: str, texts: list, name: str) -> dict:
    import tempfile
    import os
    from google import genai
    client = genai.Client(api_key=api_key)

    progress("Criando File Search Store...")
    store = client.file_search_stores.create(
        config={"display_name": _ascii_safe(f"kb-{name[:40]}")}
    )
    progress(f"Store criado: {store.name}")

    uploaded, failed = [], []
    total = len(texts)

    for i, entry in enumerate(texts, 1):
        fname = entry['name']
        text  = entry['text']
        try:
            if entry.get('already_chunked'):
                chunks = [text]
                progress(f"[{i}/{total}] {fname}: já está bem dividido, enviado como está (sem re-cortar).")
            else:
                chunks = chunk_text(text)
                if len(chunks) > 1:
                    progress(f"[{i}/{total}] {fname}: dividido em {len(chunks)} chunks.")
            stem   = Path(fname).stem
            for n, chunk in enumerate(chunks, 1):
                cname  = _ascii_filename(f"{stem}_parte{n}.txt" if len(chunks) > 1 else f"{stem}.txt")
                cbytes = chunk.encode("utf-8")
                progress(f"[{i}/{total}] Enviando {cname} ({len(cbytes)/1024:.1f} KB)...")
                tmp = tempfile.NamedTemporaryFile(suffix=".txt", delete=False)
                try:
                    tmp.write(cbytes)
                    tmp.close()
                    operation = client.file_search_stores.upload_to_file_search_store(
                        file=tmp.name,
                        file_search_store_name=store.name,
                        config={"display_name": _ascii_safe(cname)},
                    )
                finally:
                    os.unlink(tmp.name)
                progress(f"[{i}/{total}] Indexando {cname}...")
                while not operation.done:
                    time.sleep(5)
                    operation = client.operations.get(operation)
                    progress(f"[{i}/{total}] ...aguardando {cname}")
                progress(f"[{i}/{total}] {cname} indexado.")
            uploaded.append(fname)
        except Exception as e:
            failed.append(fname)
            progress(f"[{i}/{total}] ERRO em {fname}: {e}")

    if not uploaded:
        error("Nenhum texto pôde ser indexado.")
    progress("Indexacao concluida!")
    return {"store_id": store.name, "store_name": store.name, "provider": "gemini",
            "files_uploaded": uploaded, "files_failed": failed}


# ─── Add files to existing store ─────────────────────────────

def add_files_openai(api_key: str, store_id: str, file_entries: list) -> dict:
    import io
    from openai import OpenAI
    client = OpenAI(api_key=api_key)
    progress(f"Adicionando ao Vector Store {store_id}...")

    uploaded, failed = [], []
    total = len(file_entries)

    for i, entry in enumerate(file_entries, 1):
        fp   = Path(entry['path'])
        name = entry['name']
        try:
            content_bytes, upload_name = read_file_as_bytes(fp)
            text = content_bytes.decode("utf-8", errors="replace")
            chunks = chunk_text(text)
            stem = Path(upload_name).stem
            if len(chunks) > 1:
                progress(f"[{i}/{total}] {name}: dividido em {len(chunks)} chunks.")
            for n, chunk in enumerate(chunks, 1):
                chunk_name = f"{stem}_parte{n}.txt" if len(chunks) > 1 else upload_name
                chunk_bytes = chunk.encode("utf-8")
                progress(f"[{i}/{total}] Enviando {chunk_name} ({len(chunk_bytes)/1024:.1f} KB)...")
                client.vector_stores.files.upload_and_poll(
                    vector_store_id=store_id,
                    file=(chunk_name, io.BytesIO(chunk_bytes), "text/plain"),
                )
                progress(f"[{i}/{total}] {chunk_name} indexado.")
            uploaded.append(name)
        except Exception as e:
            failed.append(name)
            progress(f"[{i}/{total}] ERRO em {name}: {e}")

    if not uploaded:
        error("Nenhum arquivo pôde ser indexado.")
    progress("Adição concluída!")
    return {"store_id": store_id, "provider": "openai",
            "files_uploaded": uploaded, "files_failed": failed}


def add_files_gemini(api_key: str, store_name: str, file_entries: list) -> dict:
    import tempfile
    import os
    from google import genai

    client = genai.Client(api_key=api_key)
    progress(f"Adicionando ao File Search Store {store_name}...")

    uploaded, failed = [], []
    total = len(file_entries)

    for i, entry in enumerate(file_entries, 1):
        fp   = Path(entry['path'])
        name = entry['name']
        try:
            content_bytes, upload_name = read_file_as_bytes(fp)
            text = content_bytes.decode("utf-8", errors="replace")
            chunks = chunk_text(text)
            stem = Path(upload_name).stem

            if len(chunks) > 1:
                progress(f"[{i}/{total}] {name}: dividido em {len(chunks)} chunks.")

            chunk_ok = False
            for n, chunk in enumerate(chunks, 1):
                chunk_name = _ascii_filename(f"{stem}_parte{n}.txt" if len(chunks) > 1 else upload_name)
                chunk_bytes = chunk.encode("utf-8")
                progress(f"[{i}/{total}] Enviando {chunk_name} ({len(chunk_bytes) / 1024:.1f} KB)...")
                tmp = tempfile.NamedTemporaryFile(suffix=".txt", delete=False)
                try:
                    tmp.write(chunk_bytes)
                    tmp.close()
                    operation = client.file_search_stores.upload_to_file_search_store(
                        file=tmp.name,
                        file_search_store_name=store_name,
                        config={"display_name": _ascii_safe(chunk_name)},
                    )
                finally:
                    os.unlink(tmp.name)
                progress(f"[{i}/{total}] Indexando {chunk_name}...")
                while not operation.done:
                    time.sleep(5)
                    operation = client.operations.get(operation)
                    progress(f"[{i}/{total}] ...aguardando {chunk_name}")
                progress(f"[{i}/{total}] {chunk_name} indexado.")
                chunk_ok = True

            if chunk_ok:
                uploaded.append(name)
        except Exception as e:
            failed.append(name)
            progress(f"[{i}/{total}] ERRO em {name}: {e}")

    if not uploaded:
        error("Nenhum arquivo pôde ser indexado.")
    progress("Adição concluída!")
    return {"store_id": store_name, "provider": "gemini",
            "files_uploaded": uploaded, "files_failed": failed}


# ─── Excluir RAG (store inteiro) ──────────────────────────────

def delete_rag(provider: str, api_key: str, store_id: str) -> dict:
    if provider == "openai":
        from openai import OpenAI
        client = OpenAI(api_key=api_key)
        progress(f"Excluindo Vector Store {store_id}...")
        client.vector_stores.delete(store_id)
    else:
        from google import genai
        from google.genai import types
        client = genai.Client(api_key=api_key)
        progress(f"Excluindo File Search Store {store_id}...")
        client.file_search_stores.delete(name=store_id, config=types.DeleteFileSearchStoreConfig(force=True))

    progress("RAG excluído.")
    return {"deleted": True, "store_id": store_id}


# ─── Limpar conteúdo do RAG (mantém o store_id) ───────────────

def clear_rag(provider: str, api_key: str, store_id: str) -> dict:
    removed = 0

    if provider == "openai":
        from openai import OpenAI
        client = OpenAI(api_key=api_key)
        progress(f"Listando arquivos do Vector Store {store_id}...")
        files = list(client.vector_stores.files.list(vector_store_id=store_id))
        progress(f"{len(files)} arquivo(s) encontrado(s). Removendo...")
        for f in files:
            try:
                client.vector_stores.files.delete(f.id, vector_store_id=store_id)
                try:
                    client.files.delete(f.id)
                except Exception:
                    pass  # objeto de arquivo já pode ter sido removido
                removed += 1
                progress(f"Removido: {f.id}")
            except Exception as e:
                progress(f"Erro ao remover {f.id}: {e}")
    else:
        from google import genai
        from google.genai import types
        client = genai.Client(api_key=api_key)
        progress(f"Listando documentos do File Search Store {store_id}...")
        docs = list(client.file_search_stores.documents.list(parent=store_id))
        progress(f"{len(docs)} documento(s) encontrado(s). Removendo...")
        for d in docs:
            try:
                client.file_search_stores.documents.delete(name=d.name)
                removed += 1
                progress(f"Removido: {d.name}")
            except Exception as e:
                progress(f"Erro ao remover {d.name}: {e}")

    progress(f"Conteúdo limpo: {removed} item(ns) removido(s).")
    return {"cleared": True, "removed": removed, "store_id": store_id}


# ─── Listar stores existentes na conta do provedor ────────────

def list_stores(provider: str, api_key: str) -> dict:
    stores = []

    if provider == "openai":
        from openai import OpenAI
        client = OpenAI(api_key=api_key)
        for vs in client.vector_stores.list():
            created_iso = None
            if vs.created_at:
                created_iso = time.strftime('%Y-%m-%dT%H:%M:%SZ', time.gmtime(vs.created_at))
            stores.append({
                "store_id":   vs.id,
                "store_name": vs.name,
                "provider":   "openai",
                "file_count": vs.file_counts.total if vs.file_counts else 0,
                "createdAt":  created_iso,
            })
    else:
        from google import genai
        client = genai.Client(api_key=api_key)
        for store in client.file_search_stores.list():
            stores.append({
                "store_id":   store.name,
                "store_name": store.display_name or store.name,
                "provider":   "gemini",
                "file_count": store.active_documents_count or 0,
                "createdAt":  store.create_time.isoformat() if store.create_time else None,
            })

    return {"stores": stores}


# ─── Main ─────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser()
    sub = parser.add_subparsers(dest="cmd")

    up = sub.add_parser("upload")
    up.add_argument("--provider", required=True, choices=["openai", "gemini"])
    up.add_argument("--key", required=True)
    up.add_argument("--files-file", required=True, dest="files_file")

    qr = sub.add_parser("query")
    qr.add_argument("--provider", required=True, choices=["openai", "gemini"])
    qr.add_argument("--key", required=True)
    qr.add_argument("--store", required=True)
    qr.add_argument("--question", required=True)

    cw = sub.add_parser("crawl")
    cw.add_argument("--url",   required=True)
    cw.add_argument("--depth", type=int, default=2)  # 0 = sem limite
    cw.add_argument("--js",    action="store_true")

    sc = sub.add_parser("scrape")
    sc.add_argument("--provider",   required=True, choices=["openai", "gemini"])
    sc.add_argument("--key",        required=True)
    sc.add_argument("--urls-file",  required=True)
    sc.add_argument("--name",       required=True)

    ex = sub.add_parser("extract")
    ex.add_argument("--files-file", required=True, dest="files_file")
    ex.add_argument("--provider", choices=["openai", "gemini"], default=None)
    ex.add_argument("--key", default=None)

    ut = sub.add_parser("upload-text")
    ut.add_argument("--provider", required=True, choices=["openai", "gemini"])
    ut.add_argument("--key", required=True)
    ut.add_argument("--texts-file", required=True, dest="texts_file")
    ut.add_argument("--name", required=True)

    ad = sub.add_parser("add")
    ad.add_argument("--provider", required=True, choices=["openai", "gemini"])
    ad.add_argument("--key", required=True)
    ad.add_argument("--store", required=True)
    ad.add_argument("--files-file", required=True, dest="files_file")

    mq = sub.add_parser("multi-query")
    mq.add_argument("--provider", required=True, choices=["openai", "gemini"])
    mq.add_argument("--key", required=True)
    mq.add_argument("--stores-file", required=True, dest="stores_file")
    mq.add_argument("--question", required=True)

    ts = sub.add_parser("test")
    ts.add_argument("--provider", required=True, choices=["openai", "gemini"])
    ts.add_argument("--key", required=True)
    ts.add_argument("--store", required=True)
    ts.add_argument("--tests-file", required=True, dest="tests_file")

    dl = sub.add_parser("delete-rag")
    dl.add_argument("--provider", required=True, choices=["openai", "gemini"])
    dl.add_argument("--key", required=True)
    dl.add_argument("--store", required=True)

    cl = sub.add_parser("clear-rag")
    cl.add_argument("--provider", required=True, choices=["openai", "gemini"])
    cl.add_argument("--key", required=True)
    cl.add_argument("--store", required=True)

    ls = sub.add_parser("list-stores")
    ls.add_argument("--provider", required=True, choices=["openai", "gemini"])
    ls.add_argument("--key", required=True)

    args = parser.parse_args()

    if args.cmd == "upload":
        try:
            with open(args.files_file, encoding='utf-8') as f:
                file_entries = json.load(f)
            data = upload_openai(args.key, file_entries) if args.provider == "openai" \
                   else upload_gemini(args.key, file_entries)
            result(data)
        except Exception as e:
            error(str(e))

    elif args.cmd == "query":
        try:
            data = query_openai(args.key, args.store, args.question) \
                   if args.provider == "openai" \
                   else query_gemini(args.key, args.store, args.question)
            result(data)
        except Exception as e:
            error(str(e))

    elif args.cmd == "crawl":
        try:
            pages = crawl_site(args.url, args.depth, use_js=args.js)
            result({"pages": pages})
        except Exception as e:
            error(str(e))

    elif args.cmd == "scrape":
        try:
            with open(args.urls_file, encoding='utf-8') as f:
                urls = json.load(f)
            data = scrape_and_upload(args.key, args.provider, urls, args.name)
            result(data)
        except Exception as e:
            error(str(e))

    elif args.cmd == "extract":
        try:
            with open(args.files_file, encoding='utf-8') as f:
                file_entries = json.load(f)
            previews = []
            for entry in file_entries:
                fp = Path(entry['path'])
                name = entry['name']
                progress(f"Extraindo {name}...")
                try:
                    content_bytes, _ = read_file_as_bytes(fp)
                    text = content_bytes.decode("utf-8", errors="replace")
                    preview = {"name": name, "text": text}

                    if args.provider and args.key and fp.suffix.lower() in (".txt", ".md"):
                        progress(f"Verificando divisão em chunks de {name}...")
                        check = check_already_chunked(args.provider, args.key, text)
                        preview["already_chunked"] = check["already_chunked"]
                        preview["chunk_note"]      = check["note"]
                        progress(f"{name}: {'já dividido em chunks' if check['already_chunked'] else 'será recortado automaticamente'} — {check['note']}")

                    previews.append(preview)
                    progress(f"Extraído: {name} ({len(text)} chars)")
                except Exception as e:
                    progress(f"Erro em {name}: {e}")
            if not previews:
                error("Nenhum arquivo pôde ser extraído.")
            result({"previews": previews})
        except Exception as e:
            error(str(e))

    elif args.cmd == "upload-text":
        try:
            with open(args.texts_file, encoding='utf-8') as f:
                texts = json.load(f)  # lista de {"name": str, "text": str}
            data = _upload_texts_openai(args.key, texts, args.name) if args.provider == "openai" \
                   else _upload_texts_gemini(args.key, texts, args.name)
            result(data)
        except Exception as e:
            error(str(e))

    elif args.cmd == "add":
        try:
            with open(args.files_file, encoding='utf-8') as f:
                file_entries = json.load(f)
            data = add_files_openai(args.key, args.store, file_entries) if args.provider == "openai" \
                   else add_files_gemini(args.key, args.store, file_entries)
            result(data)
        except Exception as e:
            error(str(e))

    elif args.cmd == "multi-query":
        try:
            with open(args.stores_file, encoding='utf-8') as f:
                stores = json.load(f)  # list of {"store_id": str, "name": str, "provider": str}
            import asyncio

            async def _query_one(store_entry):
                sid = store_entry["store_id"]
                name = store_entry.get("name") or sid
                provider = store_entry.get("provider", args.provider)
                try:
                    try:
                        if provider == "openai":
                            r = await asyncio.to_thread(query_openai, args.key, sid, args.question)
                        else:
                            r = await asyncio.to_thread(query_gemini, args.key, sid, args.question)
                    except AttributeError:
                        # Python < 3.9: asyncio.to_thread not available — run sequentially
                        if provider == "openai":
                            r = query_openai(args.key, sid, args.question)
                        else:
                            r = query_gemini(args.key, sid, args.question)
                    return {"name": name, "answer": r["answer"], "citations": r.get("citations", [])}
                except Exception as e:
                    return {"name": name, "answer": f"[Erro: {e}]", "citations": []}

            async def _run_all():
                tasks = [_query_one(s) for s in stores]
                return await asyncio.gather(*tasks)

            answers = asyncio.run(_run_all())

            # Combine into a single response with source headers
            combined_parts = []
            all_citations = []
            for a in answers:
                combined_parts.append(f"### Fonte: {a['name']}\n\n{a['answer']}")
                all_citations.extend(a["citations"])

            combined = "\n\n---\n\n".join(combined_parts)
            result({"answer": combined, "citations": all_citations})
        except Exception as e:
            error(str(e))

    elif args.cmd == "test":
        try:
            with open(args.tests_file, encoding='utf-8') as f:
                tests = json.load(f)  # list of {"question": str, "expected": str}

            results_list = []
            total = len(tests)
            correct = 0

            for i, t in enumerate(tests, 1):
                q   = t.get("question", "")
                exp = t.get("expected", "")
                progress(f"[{i}/{total}] Testando: {q[:80]}...")

                try:
                    got_data = query_openai(args.key, args.store, q) \
                               if args.provider == "openai" \
                               else query_gemini(args.key, args.store, q)
                    got = got_data.get("answer", "")
                except Exception as e:
                    got = f"[Erro: {e}]"

                # LLM judge
                judge_prompt = (
                    f"Pergunta: {q}\n\n"
                    f"Resposta esperada: {exp}\n\n"
                    f"Resposta obtida: {got}\n\n"
                    "A resposta obtida responde corretamente à pergunta com base na resposta esperada? "
                    "Responda APENAS com uma das palavras: sim, parcial ou não."
                )
                verdict = "não"
                try:
                    if args.provider == "openai":
                        from openai import OpenAI
                        client = OpenAI(api_key=args.key)
                        resp = client.responses.create(model="gpt-4o-mini", input=judge_prompt)
                        verdict_raw = resp.output_text.strip().lower()
                    else:
                        from google import genai
                        client = genai.Client(api_key=args.key)
                        resp = client.models.generate_content(model="gemini-flash-latest", contents=judge_prompt)
                        verdict_raw = (resp.text or "").strip().lower()

                    if "sim" in verdict_raw:
                        verdict = "sim"
                    elif "parcial" in verdict_raw:
                        verdict = "parcial"
                    else:
                        verdict = "não"
                except Exception as e:
                    verdict = f"erro-judge: {e}"

                if verdict == "sim":
                    correct += 1
                elif verdict == "parcial":
                    correct += 0.5

                results_list.append({"question": q, "expected": exp, "got": got, "verdict": verdict})
                progress(f"[{i}/{total}] Veredicto: {verdict}")

            accuracy = round(correct / total * 100, 1) if total else 0
            progress(f"Teste concluído: {accuracy}% de acerto em {total} perguntas.")
            result({"results": results_list, "accuracy": accuracy, "total": total})
        except Exception as e:
            error(str(e))

    elif args.cmd == "delete-rag":
        try:
            data = delete_rag(args.provider, args.key, args.store)
            result(data)
        except Exception as e:
            error(str(e))

    elif args.cmd == "clear-rag":
        try:
            data = clear_rag(args.provider, args.key, args.store)
            result(data)
        except Exception as e:
            error(str(e))

    elif args.cmd == "list-stores":
        try:
            data = list_stores(args.provider, args.key)
            result(data)
        except Exception as e:
            error(str(e))

    else:
        parser.print_help()
        sys.exit(1)


if __name__ == "__main__":
    main()
